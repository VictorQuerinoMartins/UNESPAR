from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Definir cores baseadas no seu código React (Blue-900, Slate-800, etc)
    BLUE_DARK = RGBColor(30, 58, 138)  # blue-900
    BLUE_MED = RGBColor(37, 99, 235)   # blue-600
    SLATE_DARK = RGBColor(30, 41, 59)  # slate-800
    WHITE = RGBColor(255, 255, 255)

    def set_title_format(slide, text, subtitle=None):
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
        title.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            # Tenta encontrar o placeholder de subtítulo ou cria caixa de texto
            if len(slide.placeholders) > 1:
                sub = slide.placeholders[1]
                sub.text = subtitle
            else:
                # Adiciona caixa de texto para subtítulo se layout não tiver
                top = Inches(1.5)
                left = Inches(1)
                width = Inches(8.5)
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = subtitle
                tf.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139) # slate-500

    def add_notes(slide, text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = text

    # --- SLIDE 0: Capa ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Percurso Cognitivo"
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Seminário de Métodos de Avaliação de IHC\nBaseado em: Barbosa, S. D. J.; Silva, B. S. (Capítulo 12)\n\nEquipe: Carlos Arthur, Gabriel Momesso, Thiago Martins, Victor Martins"
    
    add_notes(slide, "Abertura: Apresentem a equipe. Mencione que seguirão o Cap 12 da Barbosa.")

    # --- SLIDE 1: Contextualização ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_title_format(slide, "Contextualização", "Origem e Fundamentação Teórica")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "O que é?"
    p = tf.add_paragraph()
    p.text = "Método de Inspeção de Usabilidade (sem usuários finais)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Origem: Proposto por Polson et al. (Início dos anos 90)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Foco: Avaliar sistemas 'walk up and use'."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nTeoria do Aprendizado Exploratório:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Usuários preferem aprender explorando a interface e resolvendo problemas, em vez de ler manuais."
    p.level = 1
    
    add_notes(slide, "Destaquem que é um método de INSPEÇÃO. Não tem usuário real aqui. A teoria é que as pessoas aprendem explorando, não lendo manual.")

    # --- SLIDE 2: Definição e Preparação ---
    slide = prs.slides.add_slide(slide_layout)
    set_title_format(slide, "Definição e Preparação", "Como o método é estruturado?")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "O avaliador simula o raciocínio do usuário ao realizar uma tarefa."
    
    p = tf.add_paragraph()
    p.text = "\nEntradas Necessárias (Inputs):"
    p.font.bold = True
    
    inputs = [
        ("1. Representação da Interface", "Sistema pronto, protótipo ou mockups."),
        ("2. Descrição da Tarefa", "Objetivo claro do usuário (Ex: 'Comprar um livro')."),
        ("3. Sequência de Ações", "Lista granular de passos do 'Caminho Feliz'."),
        ("4. Perfil do Usuário (Persona)", "Quem é? Qual seu conhecimento prévio?")
    ]
    
    for title, desc in inputs:
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.level = 1
        p = tf.add_paragraph()
        p.text = desc
        p.level = 2

    add_notes(slide, "Explique que nós (avaliadores) vamos 'fingir' ser o usuário. Precisamos ter a tarefa e o caminho feliz definidos antes de começar.")

    # --- SLIDE 3: As 4 Perguntas de Wharton ---
    slide = prs.slides.add_slide(slide_layout)
    set_title_format(slide, "As 4 Perguntas de Wharton", "O 'Core' da Avaliação")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Para cada ação da tarefa, faça estas 4 perguntas:"
    
    questions = [
        ("1. Intenção do Usuário", "O usuário tentará atingir o efeito correto? (Ele sabe que precisa fazer isso?)"),
        ("2. Visibilidade do Controle", "O usuário notará que a ação está disponível? (O botão é visível?)"),
        ("3. Associação (Significado)", "O usuário associará a ação com o efeito? (O ícone/texto faz sentido?)"),
        ("4. Feedback do Sistema", "O usuário perceberá o progresso? (O sistema confirma a ação?)")
    ]
    
    for q_title, q_desc in questions:
        p = tf.add_paragraph()
        p.text = "\n" + q_title
        p.font.bold = True
        p.font.color.rgb = BLUE_MED
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = q_desc
        p.level = 1

    add_notes(slide, "ESTE É O SLIDE MAIS IMPORTANTE. Leiam as 4 perguntas com calma. Elas são a base da avaliação.")

    # --- SLIDE 4: Objetivos e Aplicação ---
    slide = prs.slides.add_slide(slide_layout)
    set_title_format(slide, "Objetivos e Aplicação", "Para que serve e o que entrega?")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Principais Objetivos:"
    
    objs = [
        "Simular a resolução de problemas do novato.",
        "Identificar conflitos de modelo mental.",
        "Detectar falhas em: Terminologias, Controles escondidos, Falta de feedback."
    ]
    for obj in objs:
        p = tf.add_paragraph()
        p.text = obj
        p.level = 1
        
    p = tf.add_paragraph()
    p.text = "\nQuando utilizar?"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Fase Inicial (Design): Validação de mockups antes da programação."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Fase Final: Refinamento de interface em sistema funcional."
    p.level = 1

    add_notes(slide, "O foco é aprendizado. Serve para ver se o sistema conversa com o modelo mental do usuário.")

    # --- SLIDE 5: Exemplo Prático: Cenário ---
    slide = prs.slides.add_slide(slide_layout)
    set_title_format(slide, "Exemplo Prático: Cenário", "Aplicando o método em um caso real")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Dados do Cenário:"
    
    scenarios = [
        "Sistema: App de Delivery (Mobile).",
        "Persona: 'Sr. João', 65 anos. Usa apenas WhatsApp. Baixa familiaridade.",
        "Tarefa Global: Inserir cupom de desconto.",
        "Ação Analisada: Clicar no ícone de tíquete na tela de pagamento."
    ]
    for item in scenarios:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        
    p = tf.add_paragraph()
    p.text = "\nProblema Visual:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "O ícone de cupom (🎫) é pequeno, sem texto e está longe do botão de pagar."
    p.level = 1

    add_notes(slide, "Leiam o cenário. Destaquem que a persona é idosa, isso muda como respondemos as perguntas.")

    # --- SLIDE 6: Exemplo Prático: Análise ---
    slide = prs.slides.add_slide(slide_layout)
    set_title_format(slide, "Exemplo Prático: Análise", "Respondendo às 4 Perguntas")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Ação: 'Clicar no ícone (🎫)'"
    
    analysis = [
        ("P1. Intenção?", "SIM. Sr. João quer o desconto."),
        ("P2. Visibilidade?", "NÃO. Ícone muito pequeno e baixo contraste."),
        ("P3. Associação?", "NÃO. Sr. João procura a palavra 'Cupom', não entende o desenho."),
        ("P4. Feedback?", "N/A. O usuário falhou antes de clicar.")
    ]
    
    for q, a in analysis:
        p = tf.add_paragraph()
        p.text = q + " -> " + a
        if "NÃO" in a:
            p.font.color.rgb = RGBColor(220, 38, 38) # Red
        elif "SIM" in a:
            p.font.color.rgb = RGBColor(22, 163, 74) # Green
        p.level = 1
        
    p = tf.add_paragraph()
    p.text = "\nSolução Sugerida: Substituir ícone por texto 'Adicionar Cupom'."
    p.font.bold = True

    add_notes(slide, "Mostrem como o método encontrou o erro. O Sr. João falhou na visualização e na associação (semântica).")

    # --- SLIDE 7: Avaliação da Equipe ---
    slide = prs.slides.add_slide(slide_layout)
    set_title_format(slide, "Avaliação da Equipe", "Análise Crítica do Método")
    
    tf = slide.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "Pontos Fortes:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 163, 74) # Green
    
    p = tf.add_paragraph()
    p.text = "Custo-Benefício: Rápido e barato (sem recrutamento)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Foco no Aprendizado: Especializado em primeiro uso."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nLimitações:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(234, 88, 12) # Orange
    
    p = tf.add_paragraph()
    p.text = "Falsos Resultados: Depende da empatia do avaliador."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Escopo Limitado: Só avalia o 'caminho feliz'."
    p.level = 1

    add_notes(slide, "Avaliação crítica. É barato, mas cansativo. Precisa saber se colocar no lugar do usuário.")

    # --- SLIDE 8: Conclusão ---
    slide = prs.slides.add_slide(slide_layout)
    set_title_format(slide, "Conclusão", "Fechamento e Síntese")
    
    tf = slide.placeholders[1].text_frame
    tf.text = "'O Percurso Cognitivo não diz apenas ONDE está o problema, mas o PORQUÊ o usuário falha.'"
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "\nRecomendação:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Use em protótipos para validar terminologia e visibilidade."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nReferência Bibliográfica:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "BARBOSA, S. D. J.; SILVA, B. S. Interação Humano-Computador. Capítulo 12. Elsevier, 2010."
    p.font.italic = True
    p.level = 1

    add_notes(slide, "Fechamento. Agradecimentos e abertura para perguntas.")

    # Salvar
    prs.save('Apresentacao_IHC.pptx')
    print("Apresentação 'Apresentacao_IHC.pptx' gerada com sucesso!")

if __name__ == "__main__":
    create_presentation()